from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class ExtractedSection:
    source_type: str
    text: str
    sheet_name: str | None = None
    page_no: int | None = None
    slide_no: int | None = None
    heading: str | None = None


class UnsupportedFileTypeError(ValueError):
    pass


def extract_sections(path: Path, max_chars: int) -> list[ExtractedSection]:
    extension = path.suffix.lower()
    if extension in {".txt", ".md", ".csv", ".log"}:
        return [ExtractedSection(source_type="text", text=_read_text(path, max_chars))]
    if extension == ".docx":
        return _extract_docx(path, max_chars)
    if extension in {".xlsx", ".xlsm"}:
        return _extract_xlsx(path, max_chars)
    if extension == ".pptx":
        return _extract_pptx(path, max_chars)
    if extension == ".pdf":
        return _extract_pdf(path, max_chars)
    raise UnsupportedFileTypeError(f"unsupported file extension: {extension}")


def _read_text(path: Path, max_chars: int) -> str:
    last_error: UnicodeDecodeError | None = None
    for encoding in ("utf-8", "cp932", "shift_jis"):
        try:
            return path.read_text(encoding=encoding, errors="strict")[:max_chars]
        except UnicodeDecodeError as exc:
            last_error = exc
    if last_error is not None:
        raise last_error
    return ""


def _extract_docx(path: Path, max_chars: int) -> list[ExtractedSection]:
    from docx import Document

    document = Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return [ExtractedSection(source_type="body", text="\n".join(parts)[:max_chars])]


def _extract_xlsx(path: Path, max_chars: int) -> list[ExtractedSection]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[ExtractedSection] = []
    remaining = max_chars
    for sheet in workbook.worksheets:
        lines: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if values:
                lines.append("\t".join(values))
            if sum(len(line) for line in lines) >= remaining:
                break
        text = "\n".join(lines)[:remaining]
        if text.strip():
            sections.append(
                ExtractedSection(source_type="sheet", text=text, sheet_name=sheet.title)
            )
            remaining -= len(text)
        if remaining <= 0:
            break
    workbook.close()
    return sections


def _extract_pptx(path: Path, max_chars: int) -> list[ExtractedSection]:
    from pptx import Presentation

    presentation = Presentation(path)
    sections: list[ExtractedSection] = []
    remaining = max_chars
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)[:remaining]
        if text.strip():
            sections.append(ExtractedSection(source_type="slide", text=text, slide_no=index))
            remaining -= len(text)
        if remaining <= 0:
            break
    return sections


def _extract_pdf(path: Path, max_chars: int) -> list[ExtractedSection]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    sections: list[ExtractedSection] = []
    remaining = max_chars
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "")[:remaining]
        if text.strip():
            sections.append(ExtractedSection(source_type="page", text=text, page_no=index))
            remaining -= len(text)
        if remaining <= 0:
            break
    return sections
